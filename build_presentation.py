# -*- coding: utf-8 -*-
"""Build 8BIT-MEDIA commercial proposal PPTX for car photo & video production."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).parent
LOGO = BASE / "assets" / "logo-8bit-white.png"
OUTPUT = BASE / "KP_Image_Video_Production_8BIT_v2.pptx"

# 8BIT brand palette
BG = RGBColor(0, 0, 0)
BG2 = RGBColor(11, 11, 11)
ACCENT = RGBColor(0, 122, 255)
ACCENT2 = RGBColor(0, 85, 170)
WHITE = RGBColor(255, 255, 255)
GREY = RGBColor(160, 160, 160)
LINE = RGBColor(40, 40, 40)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.75)
MARGIN_R = Inches(0.75)
MARGIN_T = Inches(0.55)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

CURRENCY_RATE = 12200

# Split from R/A Production combined estimate ($181,828 total, 5 shoot days)
# Proportional allocation by category using photo/video quote ratios
PHOTO = {
    "items": [
        ("Подготовка / Pre-Production", 6837, "Мудборд, шот-лист, рекки локаций, PPM"),
        ("Кастинг и модели / Casting & Cast", 384, "Подбор моделей, права на использование"),
        ("Команда съёмки / Production Crew", 20475, "Продакшн, фото, свет, грим, стилист"),
        ("Оборудование / Equipment", 10753, "Камеры, оптика, свет, tether, grip"),
        ("Художественный отдел / Art Department", 5918, "Декорации, реквизит, стилизация"),
        ("Павильон / Studio", 1500, "Студийная съёмка (белый/чёрный фон)"),
        ("Локации / Location", 1339, "Аренда локаций, кейтеринг на площадке"),
        ("Транспорт / Transport", 355, "Транспорт команды и оборудования"),
        ("Постпродакшн / Post-Production", 13627, "Ретушь, цветокор, адаптация форматов"),
    ],
    "subtotal": 61188,
    "contingency": 2448,
    "production_fee": 9545,
    "before_tax": 73181,
    "banking": 3659,
    "vat": 4610,
    "total_usd": 81450,
    "shoot_days": 3,
}

VIDEO = {
    "items": [
        ("Подготовка / Pre-Production", 4823, "Сценарий, раскадровки, рекки, PPM"),
        ("Кастинг и актёры / Casting & Cast", 1636, "Кастинг, актёры, права на использование"),
        ("Команда съёмки / Production Crew", 25125, "Режиссёр, продакшн, камера, свет, звук"),
        ("Оборудование / Equipment", 10752, "ARRI/RED, оптика, свет, grip, steadicam"),
        ("Художественный отдел / Art Department", 8482, "Декорации, реквизит, костюмы"),
        ("Павильон / Studio", 1500, "Студийная съёмка (при необходимости)"),
        ("Локации / Location", 1421, "Аренда локаций, кейтеринг, пермиты"),
        ("Транспорт / Transport", 355, "Транспорт команды и оборудования"),
        ("Постпродакшн / Post-Production", 21313, "Монтаж, цвет, звук, графика, адаптация"),
    ],
    "subtotal": 75407,
    "contingency": 3016,
    "production_fee": 11763,
    "before_tax": 90186,
    "banking": 4509,
    "vat": 5683,
    "total_usd": 100378,
    "shoot_days": 2,
}

TOTAL_USD = PHOTO["total_usd"] + VIDEO["total_usd"]
TOTAL_UZS = TOTAL_USD * CURRENCY_RATE


def fmt_usd(n):
    return f"${n:,.0f}".replace(",", " ")


def fmt_uzs(n):
    s = f"{n:,.0f}".replace(",", " ")
    return f"{s} UZS"


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return box


def add_kicker(slide, text):
    add_textbox(slide, MARGIN_L, MARGIN_T, CONTENT_W, Inches(0.35), text.upper(), size=9, color=ACCENT, bold=True)


def add_title(slide, text, top=Inches(0.95)):
    add_textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.9), text, size=28, color=WHITE, bold=True)


def add_footer(slide, num, total):
    add_rect(slide, MARGIN_L, Inches(7.05), CONTENT_W, Pt(1), LINE)
    add_textbox(slide, MARGIN_L, Inches(7.12), Inches(1), Inches(0.3), f"{num:02d}", size=18, color=ACCENT, bold=True)
    add_textbox(slide, Inches(11.5), Inches(7.12), Inches(1.5), Inches(0.3), "8BIT · MEDIA", size=8, color=GREY, align=PP_ALIGN.RIGHT)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(12.35), Inches(6.95), height=Inches(0.38))


def add_logo_header(slide):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), MARGIN_L, Inches(0.35), height=Inches(0.55))


def add_bullets(slide, items, left, top, width, height, size=13, color=WHITE, spacing=Pt(6)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"—  {item}"
        p.space_after = spacing
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def add_budget_table(slide, block, title, left=Inches(0.75), top=Inches(1.85)):
    summary = [
        ("Себестоимость / Sub Total", block["subtotal"], ""),
        ("Непредвиденные 4% / Contingency", block["contingency"], ""),
        ("Комиссия продакшена 15%", block["production_fee"], ""),
        ("Банковские расходы 5%", block["banking"], ""),
        ("НДС 6% / VAT", block["vat"], ""),
        ("ИТОГО / TOTAL", block["total_usd"], fmt_uzs(block["total_usd"] * CURRENCY_RATE)),
    ]
    rows = 1 + len(block["items"]) + len(summary)
    tbl_shape = slide.shapes.add_table(rows, 3, left, top, Inches(11.85), Inches(0.30 * rows))
    table = tbl_shape.table
    table.columns[0].width = Inches(4.2)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(6.15)

    headers = ["Статья расходов", "USD", "Комментарий"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT2
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(10)
                r.font.color.rgb = WHITE

    for i, (name, usd, note) in enumerate(block["items"], start=1):
        table.cell(i, 0).text = name
        table.cell(i, 1).text = fmt_usd(usd) if usd else "—"
        table.cell(i, 2).text = note
        for c in range(3):
            cell = table.cell(i, c)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG2
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
                    r.font.color.rgb = WHITE if c != 2 else GREY

    r = len(block["items"]) + 1
    for j, (label, val, note) in enumerate(summary):
        idx = r + j
        table.cell(idx, 0).text = label
        table.cell(idx, 1).text = fmt_usd(val)
        table.cell(idx, 2).text = note
        is_total = j == len(summary) - 1
        for c in range(3):
            cell = table.cell(idx, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT if is_total else BG2
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(10 if is_total else 9)
                    run.font.bold = is_total
                    run.font.color.rgb = WHITE


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    slide_num = 0
    total_slides = 11

    # --- Slide 1: Cover ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_rect(s, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), MARGIN_L, Inches(1.2), height=Inches(0.85))
    add_textbox(s, MARGIN_L, Inches(2.35), CONTENT_W, Inches(1.5),
                "КОММЕРЧЕСКОЕ\nПРЕДЛОЖЕНИЕ", size=40, color=WHITE, bold=True)
    add_textbox(s, MARGIN_L, Inches(4.0), CONTENT_W, Inches(0.6),
                "Фото- и видеопроизводство для launch-кампании новой модели", size=18, color=GREY)
    add_textbox(s, MARGIN_L, Inches(4.75), CONTENT_W, Inches(0.5),
                "Image & Video Production · 90 фото + TVC и digital-ролики", size=13, color=ACCENT)
    add_rect(s, MARGIN_L, Inches(5.55), Inches(2.2), Pt(2), ACCENT)
    add_textbox(s, MARGIN_L, Inches(5.75), Inches(5), Inches(0.35),
                "Digital-агентство полного цикла · Ташкент", size=11, color=GREY)
    add_textbox(s, MARGIN_L, Inches(6.15), Inches(5), Inches(0.35),
                "8bit.uz  ·  info@8bit.uz", size=11, color=GREY)
    add_footer(s, slide_num, total_slides)

    # --- Slide 2: About ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_logo_header(s)
    add_kicker(s, "О компании")
    add_title(s, "8BIT-MEDIA")
    add_textbox(s, MARGIN_L, Inches(1.85), Inches(5.8), Inches(1.2),
                "Digital-агентство полного цикла с 10+ годами опыта. "
                "Мы управляем проектами от стратегии и креатива до производства и медиаразмещения.",
                size=14, color=GREY)
    cards = [
        ("Продакшн", "Опыт автомобильных брендов: KIA, Geely, Hongqi, Zeekr. Студия с белым/чёрным фоном."),
        ("Креатив", "Разработка концепций, KV, сценариев и визуальной айдентики кампании"),
        ("Медиа", "Адаптация материалов для digital, SMM, официального сайта и POS"),
    ]
    for i, (t, d) in enumerate(cards):
        x = MARGIN_L + Inches(i * 3.95)
        add_rect(s, x, Inches(3.2), Inches(3.7), Inches(2.5), BG2, LINE)
        add_textbox(s, x + Inches(0.2), Inches(3.4), Inches(3.3), Inches(0.4), t, size=14, color=ACCENT, bold=True)
        add_textbox(s, x + Inches(0.2), Inches(3.85), Inches(3.3), Inches(1.6), d, size=11, color=GREY)
    add_footer(s, slide_num, total_slides)

    # --- Slide 3: Project ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_logo_header(s)
    add_kicker(s, "Проект")
    add_title(s, "Задача и контекст")
    add_textbox(s, MARGIN_L, Inches(1.85), CONTENT_W, Inches(1.0),
                "Для launch event, медиакампании, официального сайта, социальных сетей и "
                "POS-материалов требуется профессиональное фото- и видеопроизводство нового автомобиля. "
                "8BIT-MEDIA выступает генеральным подрядчиком и координирует продакшн.",
                size=13, color=WHITE)
    cols = [
        ("Клиент / проект", "Launch новой модели · Automotive"),
        ("Локация", "Ташкент, Узбекистан"),
        ("Срок сдачи", "01.09.2026"),
        ("Съёмочные дни", "5 дней (фото 3 + видео 2)"),
        ("Форматы фото", "RAW · TIFF · JPG (10 000×7 000 px)"),
        ("Форматы видео", "4K UHD · 25/50 fps · LOG"),
        ("Курс USD/UZS", f"1 USD = {CURRENCY_RATE:,} UZS".replace(",", " ")),
        ("Продакшн-партнёр", "R/A Production (смета субподрядчика)"),
    ]
    for i, (k, v) in enumerate(cols):
        col = i % 2
        row = i // 2
        x = MARGIN_L + Inches(col * 6.0)
        y = Inches(3.0 + row * 1.05)
        add_textbox(s, x, y, Inches(2.2), Inches(0.3), k, size=10, color=ACCENT, bold=True)
        add_textbox(s, x + Inches(2.2), y, Inches(3.5), Inches(0.35), v, size=12, color=WHITE)
    add_footer(s, slide_num, total_slides)

    # --- Slide 4: Photo scope ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_logo_header(s)
    add_kicker(s, "Объём работ · Фото")
    add_title(s, "Фотопроизводство — 90 кадров")
    photo_items = [
        "Экстерьер / Exterior — 20 шт. (от 10 000 × 7 000 px)",
        "Интерьер / Interior — 20 шт.",
        "Демонстрация функций — 15 шт.",
        "Разные локации и сцены — 25 шт.",
        "KV / Key Visual — 10 шт.",
    ]
    add_bullets(s, photo_items, MARGIN_L, Inches(2.0), Inches(5.8), Inches(3.5), size=13)
    add_rect(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.2), BG2, LINE)
    add_textbox(s, Inches(7.25), Inches(2.2), Inches(5.1), Inches(0.35), "Постобработка", size=12, color=ACCENT, bold=True)
    post = [
        "Цветокоррекция и ретушь",
        "Удаление дефектов, отражений, посторонних объектов",
        "Оптимизация фона и лакокрасочного покрытия",
        "Адаптация под форматы digital и print",
    ]
    add_bullets(s, post, Inches(7.25), Inches(2.65), Inches(5.0), Inches(2.5), size=11, color=GREY)
    add_textbox(s, Inches(7.25), Inches(5.0), Inches(5.0), Inches(0.8),
                "Поставка: RAW + TIFF + JPG (финальные ретуши)", size=11, color=WHITE, bold=True)
    add_textbox(s, Inches(7.25), Inches(5.55), Inches(5.0), Inches(1.2),
                "QC: без отражений оборудования, теней команды, лишнего реквизита и фонового мусора",
                size=10, color=GREY)
    add_footer(s, slide_num, total_slides)

    # --- Slide 5: Video scope ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_logo_header(s)
    add_kicker(s, "Объём работ · Видео")
    add_title(s, "Видеопроизводство")
    vid_left = [
        "TVC / главный рекламный ролик — 2–3 версии (30″ / 15″)",
        "Экстерьер — 2–3 ролика",
        "Интерьер — 2–3 ролика",
        "Функциональные видео — 2–3 ролика",
    ]
    add_bullets(s, vid_left, MARGIN_L, Inches(2.0), Inches(5.8), Inches(2.8), size=13)
    add_rect(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.5), BG2, LINE)
    add_textbox(s, Inches(7.25), Inches(2.2), Inches(5.1), Inches(0.35), "Технические требования", size=12, color=ACCENT, bold=True)
    tech = [
        "4K UHD (3840×2160), 25 или 50 fps",
        "Сценарий, монтаж, цветокор, графика",
        "Языки: узбекский / русский",
        "TVC 30″/15″ + digital-ролики (экстерьер, интерьер, функции)",
        "Поставка: проектные файлы, LOG, версии без титров, EN/RU/UZ субтитры",
        "QC: без тряски, шумов, киноляпов; единый стиль; синхрон ≤ 1 кадр",
    ]
    add_bullets(s, tech, Inches(7.25), Inches(2.65), Inches(5.0), Inches(3.5), size=11, color=GREY)
    add_footer(s, slide_num, total_slides)

    # --- Slide 6: Approach ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_logo_header(s)
    add_kicker(s, "Подход")
    add_title(s, "Этапы реализации")
    phases = [
        ("01", "Pre-Production", "Креатив, moodboard, шот-лист / сценарий, рекки локаций, PPM с клиентом"),
        ("02", "Production", "Студийная и location-съёмка, отдельные команды под фото и видео"),
        ("03", "Post-Production", "Ретушь / монтаж, цвет, звук, графика, адаптация форматов"),
        ("04", "Delivery", "Передача исходников и финальных материалов, правки до утверждения"),
    ]
    for i, (num, title, desc) in enumerate(phases):
        x = MARGIN_L + Inches(i * 3.05)
        add_rect(s, x, Inches(2.0), Inches(2.85), Inches(3.8), BG2, LINE)
        add_textbox(s, x + Inches(0.2), Inches(2.15), Inches(1), Inches(0.5), num, size=22, color=ACCENT, bold=True)
        add_textbox(s, x + Inches(0.2), Inches(2.75), Inches(2.5), Inches(0.5), title, size=13, color=WHITE, bold=True)
        add_textbox(s, x + Inches(0.2), Inches(3.35), Inches(2.5), Inches(2.2), desc, size=10, color=GREY)
    add_textbox(s, MARGIN_L, Inches(6.1), CONTENT_W, Inches(0.4),
                f"Съёмочные дни (оценка R/A Production): фото — {PHOTO['shoot_days']} дн. · видео — {VIDEO['shoot_days']} дн. · всего 5 дн.",
                size=11, color=ACCENT)
    add_footer(s, slide_num, total_slides)

    # --- Slide 7: Photo budget ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_logo_header(s)
    add_kicker(s, "Смета · Фото")
    add_title(s, f"Фотопроизводство — {fmt_usd(PHOTO['total_usd'])}")
    add_budget_table(s, PHOTO, "Photo")
    add_footer(s, slide_num, total_slides)

    # --- Slide 8: Video budget ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_logo_header(s)
    add_kicker(s, "Смета · Видео")
    add_title(s, f"Видеопроизводство — {fmt_usd(VIDEO['total_usd'])}")
    add_budget_table(s, VIDEO, "Video")
    add_footer(s, slide_num, total_slides)

    # --- Slide 9: Summary ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_logo_header(s)
    add_kicker(s, "Итого")
    add_title(s, "Сводная стоимость проекта")
    summary_rows = [
        ("Фотопроизводство", PHOTO["total_usd"]),
        ("Видеопроизводство", VIDEO["total_usd"]),
        ("ИТОГО", TOTAL_USD),
    ]
    y = Inches(2.1)
    for i, (label, val) in enumerate(summary_rows):
        h = Inches(0.85)
        fill = ACCENT if i == 2 else BG2
        add_rect(s, MARGIN_L, y, Inches(8.5), h, fill, LINE if i < 2 else None)
        add_textbox(s, MARGIN_L + Inches(0.25), y + Inches(0.22), Inches(5), Inches(0.4),
                    label, size=16 if i == 2 else 14, color=WHITE, bold=(i == 2))
        add_textbox(s, Inches(6.5), y + Inches(0.18), Inches(2.5), Inches(0.45),
                    fmt_usd(val), size=20 if i == 2 else 16, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
        y += h + Inches(0.12)
    add_textbox(s, MARGIN_L, Inches(4.15), Inches(8.5), Inches(0.4),
                fmt_uzs(TOTAL_UZS), size=16, color=ACCENT, bold=True)
    notes = [
        "Смета на основе котировки субподрядчика R/A Production (общий бюджет $181 828)",
        "Разделение фото/видео — пропорциональное по статьям расходов, без детализации персонала",
        "Включены: contingency 4%, production fee 15%, банк. расходы 5%, НДС 6%",
        "Не включены: права на музыку beyond stock, международные перелёты, доп. съёмочные дни",
        "Финальная стоимость фиксируется после утверждения креатива и локаций",
    ]
    add_bullets(s, notes, MARGIN_L, Inches(4.75), Inches(10), Inches(2.0), size=10, color=GREY)
    add_footer(s, slide_num, total_slides)

    # --- Slide 10: Terms ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_logo_header(s)
    add_kicker(s, "Условия")
    add_title(s, "Коммерческие условия")
    terms = [
        "Предоплата 50% — при подписании договора и старте pre-production",
        "40% — перед началом съёмочного периода",
        "10% — после финальной сдачи и принятия материалов",
        "Срок действия предложения: 30 календарных дней",
        "Правки: до полного утверждения в рамках согласованного scope",
        "Права на использование: согласно договору и территории кампании",
        "8BIT-MEDIA выступает генеральным подрядчиком и координирует продакшн",
    ]
    add_bullets(s, terms, MARGIN_L, Inches(2.0), Inches(10.5), Inches(4.5), size=13)
    add_footer(s, slide_num, total_slides)

    # --- Slide 11: Contact ---
    slide_num += 1
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BG2)
    add_rect(s, Inches(0), Inches(0), Inches(0.12), SLIDE_H, ACCENT)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), MARGIN_L, Inches(2.0), height=Inches(0.9))
    add_textbox(s, MARGIN_L, Inches(3.2), CONTENT_W, Inches(0.8),
                "Готовы обсудить детали\nи запустить проект", size=32, color=WHITE, bold=True)
    add_textbox(s, MARGIN_L, Inches(4.35), Inches(5), Inches(0.35),
                "8bit.uz  ·  info@8bit.uz", size=14, color=ACCENT)
    add_textbox(s, MARGIN_L, Inches(4.85), Inches(6), Inches(0.35),
                "Digital-агентство полного цикла · Ташкент", size=12, color=GREY)
    add_textbox(s, MARGIN_L, Inches(5.5), Inches(6), Inches(0.35),
                "Спасибо за внимание", size=11, color=GREY)
    add_footer(s, slide_num, total_slides)

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Total USD: {fmt_usd(TOTAL_USD)}")
    print(f"Total UZS: {fmt_uzs(TOTAL_UZS)}")


if __name__ == "__main__":
    build()
