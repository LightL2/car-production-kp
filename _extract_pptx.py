# -*- coding: utf-8 -*-
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).parent
PPTX = ROOT / "Anor Bureau x 8BIT x BYD Presentation.pptx"
OUT = ROOT / "pptx_extract"
OUT.mkdir(exist_ok=True)

prs = Presentation(str(PPTX))
slides_data = []

for i, slide in enumerate(prs.slides, 1):
    slide_info = {"num": i, "shapes": []}
    for j, shape in enumerate(slide.shapes):
        item = {"idx": j, "type": str(shape.shape_type)}
        if shape.has_text_frame:
            text = shape.text.strip()
            if text:
                item["text"] = text
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            fn = OUT / f"slide{i:02d}_img{j:02d}.{img.ext}"
            fn.write_bytes(img.blob)
            item["file"] = fn.name
            item["w"] = shape.width
            item["h"] = shape.height
        if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            item["media"] = True
        slide_info["shapes"].append(item)
    slides_data.append(slide_info)

(ROOT / "pptx_slides.json").write_text(
    json.dumps(slides_data, ensure_ascii=False, indent=2), encoding="utf-8"
)
print(json.dumps(slides_data, ensure_ascii=False, indent=2))
